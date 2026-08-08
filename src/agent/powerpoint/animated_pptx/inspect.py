"""Read-only inspection of arbitrary PowerPoint packages.

Inspection exposes native OOXML identities and conservative edit capabilities
without saving, rewriting, or otherwise mutating the source presentation.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from .animator import NS
from .validator import validate_package


def _inches(value: int | None) -> float | None:
    return None if value is None else round(Emu(value).inches, 6)


def _enum_name(value: Any) -> str:
    return getattr(value, "name", str(value))


def _shape_capabilities(shape) -> dict[str, bool]:
    shape_type = _enum_name(shape.shape_type)
    return {
        "select_by_shape_id": True,
        "update_geometry": hasattr(shape, "left") and hasattr(shape, "top"),
        "update_text": bool(getattr(shape, "has_text_frame", False)),
        "update_fill": shape_type in {"AUTO_SHAPE", "FREEFORM", "PLACEHOLDER", "TEXT_BOX"},
        "replace": True,
        "delete": True,
        "inspect_children": shape_type == "GROUP",
    }


def _placeholder(shape) -> dict[str, Any] | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    placeholder = shape.placeholder_format
    return {
        "idx": placeholder.idx,
        "type": _enum_name(placeholder.type),
    }


def _shape_warning(shape, slide_id: int) -> dict[str, Any] | None:
    shape_type = _enum_name(shape.shape_type)
    opaque = {
        "CHART", "DIAGRAM", "EMBEDDED_OLE_OBJECT", "FORM_CONTROL",
        "IGX_GRAPHIC", "INK", "LINKED_OLE_OBJECT", "MEDIA", "WEB_VIDEO",
    }
    if shape_type not in opaque:
        return None
    return {
        "code": f"OPAQUE_{shape_type}",
        "message": f"{shape_type} is inspectable only at container level",
        "slide_id": slide_id,
        "shape_id": shape.shape_id,
        "preserved_by_package_only_edits": True,
        "high_level_edit_safe": False,
    }


def _inspect_shape(shape, slide_id: int, z_order: int) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    shape_type = _enum_name(shape.shape_type)
    children: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    if shape_type == "GROUP":
        for child_index, child in enumerate(shape.shapes):
            inspected, child_warnings = _inspect_shape(child, slide_id, child_index)
            children.append(inspected)
            warnings.extend(child_warnings)
    warning = _shape_warning(shape, slide_id)
    if warning:
        warnings.append(warning)
    text = None
    if getattr(shape, "has_text_frame", False):
        text = shape.text
        if len(text) > 500:
            text = text[:497] + "..."
    result = {
        "selector": {"slide_id": slide_id, "shape_id": shape.shape_id},
        "shape_id": shape.shape_id,
        "name": shape.name,
        "type": shape_type,
        "z_order": z_order,
        "bounds": {
            "x": _inches(getattr(shape, "left", None)),
            "y": _inches(getattr(shape, "top", None)),
            "width": _inches(getattr(shape, "width", None)),
            "height": _inches(getattr(shape, "height", None)),
        },
        "rotation": getattr(shape, "rotation", 0.0),
        "text": text,
        "placeholder": _placeholder(shape),
        "capabilities": _shape_capabilities(shape),
        "children": children,
    }
    return result, warnings


def _animations(slide) -> list[dict[str, Any]]:
    root = slide._element
    animations = []
    for order, target in enumerate(root.findall(".//p:timing//p:spTgt", NS)):
        shape_id = target.get("spid")
        behavior = next(
            (
                ancestor for ancestor in target.iterancestors()
                if etree.QName(ancestor).localname in {
                    "anim", "animClr", "animEffect", "animMotion", "animRot",
                    "animScale", "audio", "cmd", "set", "video",
                }
            ),
            None,
        )
        wrapper = next(
            (
                ancestor for ancestor in target.iterancestors()
                if ancestor.tag == f"{{{NS['p']}}}cTn"
                and ancestor.get("presetClass") is not None
            ),
            None,
        )
        animations.append({
            "order": order,
            "target_shape_id": int(shape_id) if shape_id and shape_id.isdigit() else shape_id,
            "behavior": etree.QName(behavior).localname if behavior is not None else "unknown",
            "preset_class": wrapper.get("presetClass") if wrapper is not None else None,
            "preset_id": wrapper.get("presetID") if wrapper is not None else None,
            "preset_subtype": wrapper.get("presetSubtype") if wrapper is not None else None,
        })
    return animations


def _transition(slide) -> dict[str, Any] | None:
    transition = slide._element.find("p:transition", NS)
    if transition is None:
        return None
    child = next(iter(transition), None)
    return {
        "type": etree.QName(child).localname if child is not None else "unknown",
        "speed": transition.get("spd"),
        "attributes": dict(child.attrib) if child is not None else {},
    }


def _relationships(part) -> list[dict[str, Any]]:
    relationships = []
    for relationship in part.rels.values():
        external = relationship.is_external
        relationships.append({
            "id": relationship.rId,
            "type": relationship.reltype,
            "target_mode": "External" if external else "Internal",
            "target": relationship.target_ref,
            "target_part": None if external else str(relationship.target_part.partname).lstrip("/"),
        })
    return sorted(relationships, key=lambda item: item["id"])


def inspect_presentation(path: str | Path) -> dict[str, Any]:
    """Return a stable, JSON-serializable description without changing the file."""
    source = Path(path)
    if not source.is_file():
        raise FileNotFoundError(f"input presentation does not exist: {source}")
    if source.suffix.lower() != ".pptx":
        raise ValueError("input_path must have a .pptx extension")
    validate_package(source)
    presentation = Presentation(str(source))
    slides = []
    warnings: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides):
        slide_id_entry = presentation.slides._sldIdLst[index]
        slide_id = int(slide_id_entry.id)
        shapes = []
        for z_order, shape in enumerate(slide.shapes):
            inspected, shape_warnings = _inspect_shape(shape, slide_id, z_order)
            shapes.append(inspected)
            warnings.extend(shape_warnings)
        name_counts: dict[str, int] = {}
        for shape in slide.shapes:
            name_counts[shape.name] = name_counts.get(shape.name, 0) + 1
        for name, count in name_counts.items():
            if count > 1:
                warnings.append({
                    "code": "AMBIGUOUS_SHAPE_NAME",
                    "message": f"shape name {name!r} appears {count} times",
                    "slide_id": slide_id,
                    "preserved_by_package_only_edits": True,
                    "high_level_edit_safe": False,
                })
        slides.append({
            "selector": {"slide_id": slide_id},
            "slide_id": slide_id,
            "index": index,
            "part": str(slide.part.partname).lstrip("/"),
            "layout": {
                "name": slide.slide_layout.name,
                "part": str(slide.slide_layout.part.partname).lstrip("/"),
            },
            "shapes": shapes,
            "animations": _animations(slide),
            "transition": _transition(slide),
            "relationships": _relationships(slide.part),
            "capabilities": {
                "select_by_slide_id": True,
                "high_level_edit": True,
                "package_edit": True,
            },
        })
    properties = presentation.core_properties
    return {
        "operation": "inspect",
        "path": str(source),
        "presentation": {
            "slide_count": len(slides),
            "width": _inches(presentation.slide_width),
            "height": _inches(presentation.slide_height),
            "properties": {
                "title": properties.title,
                "subject": properties.subject,
                "author": properties.author,
                "keywords": properties.keywords,
                "comments": properties.comments,
            },
        },
        "slides": slides,
        "warnings": warnings,
        "capabilities": {
            "read_only": True,
            "native_slide_selectors": True,
            "native_shape_selectors": True,
            "stable_selector_edits": True,
            "package_edit": True,
        },
        "preservation": {
            "source_mutated": False,
            "inspection_rewrites_package": False,
            "opaque_objects": "reported_and_untouched",
        },
    }
