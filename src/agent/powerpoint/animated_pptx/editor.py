"""Atomic editing operations for existing PowerPoint presentations."""

from __future__ import annotations

from copy import deepcopy
import os
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .animator import apply_slide_effects, update_timing_targets
from .builder import add_element
from .model import Animation, Element, Slide
from .validator import validate_package


def _slide_at(presentation: Presentation, index: int, context: str = "slide_index"):
    if isinstance(index, bool) or not isinstance(index, int):
        raise TypeError(f"{context} must be an integer")
    if not 0 <= index < len(presentation.slides):
        raise IndexError(f"{context} {index} is out of range")
    return presentation.slides[index]


def _native_id(value: Any, context: str) -> int:
    if isinstance(value, bool) or not isinstance(value, int) or value < 1:
        raise TypeError(f"{context} must be a positive integer")
    return value


def _slide_by_id(presentation: Presentation, slide_id: int):
    native_id = _native_id(slide_id, "selector.slide_id")
    for index, entry in enumerate(presentation.slides._sldIdLst):
        if int(entry.id) == native_id:
            return index, presentation.slides[index]
    raise ValueError(f"slide_id {native_id} does not exist in the presentation")


def _selector(raw: Any, context: str, *, shape: bool) -> dict[str, int]:
    if not isinstance(raw, dict):
        raise TypeError(f"{context} must be an object")
    allowed = {"slide_id", "shape_id"} if shape else {"slide_id"}
    unknown = set(raw) - allowed
    if unknown:
        raise ValueError(f"{context} contains unsupported field(s): {', '.join(sorted(unknown))}")
    if "slide_id" not in raw or (shape and "shape_id" not in raw):
        required = "slide_id and shape_id" if shape else "slide_id"
        raise ValueError(f"{context} must contain {required}")
    parsed = {"slide_id": _native_id(raw["slide_id"], f"{context}.slide_id")}
    if shape:
        parsed["shape_id"] = _native_id(raw["shape_id"], f"{context}.shape_id")
    return parsed


def _selected_slide(
    presentation: Presentation,
    modifier: dict[str, Any],
    *,
    index_key: str = "slide_index",
    selector_key: str = "selector",
):
    has_index = index_key in modifier
    has_selector = selector_key in modifier
    if has_index == has_selector:
        raise ValueError(f"provide exactly one of {index_key} or {selector_key}")
    if has_selector:
        parsed = _selector(modifier[selector_key], selector_key, shape=False)
        return _slide_by_id(presentation, parsed["slide_id"])
    index = modifier[index_key]
    return index, _slide_at(presentation, index, index_key)


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _shape_by_id(slide, shape_id: int):
    native_id = _native_id(shape_id, "selector.shape_id")
    matches = [shape for shape in _iter_shapes(slide.shapes) if shape.shape_id == native_id]
    if not matches:
        raise ValueError(f"shape_id {native_id} does not exist on selected slide")
    if len(matches) > 1:
        raise ValueError(f"shape_id {native_id} is ambiguous on selected slide")
    return matches[0]


def _selected_shape(presentation: Presentation, modifier: dict[str, Any]):
    has_selector = "selector" in modifier
    has_legacy = "slide_index" in modifier or "element_id" in modifier
    if has_selector and has_legacy:
        raise ValueError("provide selector or slide_index/element_id, not both")
    if has_selector:
        parsed = _selector(modifier["selector"], "selector", shape=True)
        index, slide = _slide_by_id(presentation, parsed["slide_id"])
        return index, slide, _shape_by_id(slide, parsed["shape_id"])
    if "slide_index" not in modifier or "element_id" not in modifier:
        raise ValueError("provide selector or both slide_index and element_id")
    index = modifier["slide_index"]
    slide = _slide_at(presentation, index)
    return index, slide, _shape(slide, modifier["element_id"])


def _insert_last_slide_at(presentation: Presentation, index: int) -> None:
    if not 0 <= index <= len(presentation.slides) - 1:
        raise IndexError(f"insert index {index} is out of range")
    slide_ids = presentation.slides._sldIdLst
    slide_id = slide_ids[-1]
    slide_ids.remove(slide_id)
    slide_ids.insert(index, slide_id)


def _add_slide(presentation: Presentation, spec: Slide, index: int | None = None) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    mapping: dict[str, int] = {}
    for element in spec.elements:
        shape = add_element(slide, element)
        mapping[element.id] = shape.shape_id
    if index is not None:
        _insert_last_slide_at(presentation, index)
        slide = presentation.slides[index]
    apply_slide_effects(slide, spec, mapping, (index or len(presentation.slides) - 1) + 1)


def _clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def _shape_names(slide) -> set[str]:
    return {shape.name for shape in _iter_shapes(slide.shapes)}


def _mapping(slide) -> dict[str, int]:
    mapping: dict[str, int] = {}
    for shape in _iter_shapes(slide.shapes):
        if shape.name in mapping:
            raise ValueError(
                f"element name {shape.name!r} is ambiguous on slide; use native selectors"
            )
        mapping[shape.name] = shape.shape_id
    return mapping


def _shape(slide, element_id: str):
    if not isinstance(element_id, str) or not element_id:
        raise TypeError("element_id must be a non-empty string")
    matches = [shape for shape in _iter_shapes(slide.shapes) if shape.name == element_id]
    if not matches:
        raise ValueError(f"element {element_id!r} does not exist on slide")
    if len(matches) > 1:
        raise ValueError(
            f"element name {element_id!r} is ambiguous on slide; use a native selector"
        )
    return matches[0]


def _shape_is_top_level(slide, shape) -> bool:
    return any(candidate._element is shape._element for candidate in slide.shapes)


def _require_top_level(slide, shape, operation: str) -> None:
    if not _shape_is_top_level(slide, shape):
        raise ValueError(f"{operation} currently requires a top-level shape selector")


def _set_text_style(shape, changes: dict[str, Any]) -> None:
    if "text" in changes:
        if not getattr(shape, "has_text_frame", False):
            raise ValueError("element has no editable text frame")
        if not isinstance(changes["text"], str):
            raise TypeError("text must be a string")
        shape.text = changes["text"]
    if not getattr(shape, "has_text_frame", False):
        return
    font_size = changes.get("font_size")
    text_color = changes.get("text_color")
    if font_size is not None and (isinstance(font_size, bool) or not isinstance(font_size, (int, float)) or font_size <= 0):
        raise ValueError("font_size must be greater than zero")
    if text_color is not None:
        if not isinstance(text_color, str) or len(text_color) != 6:
            raise ValueError("text_color must be a six-digit RGB hex string")
        RGBColor.from_string(text_color)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if font_size is not None:
                run.font.size = Pt(font_size)
            if text_color is not None:
                run.font.color.rgb = RGBColor.from_string(text_color)


def _update_element(shape, changes: dict[str, Any]) -> None:
    if not isinstance(changes, dict):
        raise TypeError("changes must be an object")
    for key in ("x", "y", "width", "height", "rotation"):
        if key in changes:
            value = changes[key]
            if isinstance(value, bool) or not isinstance(value, (int, float)):
                raise TypeError(f"{key} must be a number")
            if key in ("width", "height") and value <= 0:
                raise ValueError(f"{key} must be greater than zero")
            if key in ("x", "y") and value < 0:
                raise ValueError(f"{key} must not be negative")
            setattr(shape, {"x": "left", "y": "top"}.get(key, key), Inches(value) if key != "rotation" else value)
    transform = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
    for key, attribute in (("flip_horizontal", "flipH"), ("flip_vertical", "flipV")):
        if key in changes:
            if not isinstance(changes[key], bool):
                raise TypeError(f"{key} must be a boolean")
            if transform is None:
                raise ValueError("element has no editable transform")
            if changes[key]:
                transform.set(attribute, "1")
            else:
                transform.attrib.pop(attribute, None)
    if "fill_color" in changes:
        color = changes["fill_color"]
        if not isinstance(color, str) or len(color) != 6:
            raise ValueError("fill_color must be a six-digit RGB hex string")
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
    _set_text_style(shape, changes)


def _element_ids(modifier: dict[str, Any], *, minimum: int = 1) -> list[str]:
    ids = modifier.get("element_ids")
    if not isinstance(ids, list) or len(ids) < minimum or not all(
        isinstance(item, str) and item for item in ids
    ):
        raise TypeError(f"element_ids must contain at least {minimum} non-empty string(s)")
    if len(ids) != len(set(ids)):
        raise ValueError("element_ids must be unique")
    return ids


def _selected_shapes(
    presentation: Presentation,
    modifier: dict[str, Any],
    *,
    minimum: int = 1,
):
    has_selectors = "selectors" in modifier
    has_legacy = "slide_index" in modifier or "element_ids" in modifier
    if has_selectors and has_legacy:
        raise ValueError("provide selectors or slide_index/element_ids, not both")
    if has_selectors:
        raw = modifier["selectors"]
        if not isinstance(raw, list) or len(raw) < minimum:
            raise TypeError(f"selectors must contain at least {minimum} selector object(s)")
        parsed = [_selector(item, f"selectors[{index}]", shape=True) for index, item in enumerate(raw)]
        identities = [(item["slide_id"], item["shape_id"]) for item in parsed]
        if len(identities) != len(set(identities)):
            raise ValueError("selectors must be unique")
        slide_ids = {item["slide_id"] for item in parsed}
        if len(slide_ids) != 1:
            raise ValueError("all shape selectors must target the same slide")
        index, slide = _slide_by_id(presentation, parsed[0]["slide_id"])
        return index, slide, [_shape_by_id(slide, item["shape_id"]) for item in parsed]
    if "slide_index" not in modifier or "element_ids" not in modifier:
        raise ValueError("provide selectors or both slide_index and element_ids")
    index = modifier["slide_index"]
    slide = _slide_at(presentation, index)
    return index, slide, [_shape(slide, item) for item in _element_ids(modifier, minimum=minimum)]


def _duplicate_elements(slide, shapes, modifier: dict[str, Any]) -> None:
    new_ids = modifier.get("new_ids")
    if not isinstance(new_ids, list) or len(new_ids) != len(shapes) or not all(
        isinstance(item, str) and item for item in new_ids
    ):
        raise TypeError("new_ids must contain one non-empty string for each selected element")
    existing = _shape_names(slide)
    if len(new_ids) != len(set(new_ids)) or any(item in existing for item in new_ids):
        raise ValueError("new_ids must be unique and unused on the slide")
    dx, dy = modifier.get("offset_x", 0.2), modifier.get("offset_y", 0.2)
    if any(isinstance(value, bool) or not isinstance(value, (int, float)) for value in (dx, dy)):
        raise TypeError("offset_x and offset_y must be numbers")
    next_shape_id = max(shape.shape_id for shape in _iter_shapes(slide.shapes)) + 1
    for source, new_id in zip(shapes, new_ids):
        _require_top_level(slide, source, "duplicate_elements")
        copied = deepcopy(source._element)
        properties = copied.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
        if properties is None:
            raise ValueError(f"selected shape {source.shape_id} has no non-visual properties")
        properties.set("id", str(next_shape_id))
        properties.set("name", new_id)
        next_shape_id += 1
        slide.shapes._spTree.insert_element_before(copied, "p:extLst")
        duplicate = _shape(slide, new_id)
        duplicate.left += Inches(dx)
        duplicate.top += Inches(dy)


def _reorder_elements(slide, shapes, modifier: dict[str, Any]) -> None:
    position = modifier.get("position")
    tree = slide.shapes._spTree
    for shape in shapes:
        _require_top_level(slide, shape, "reorder_elements")
    elements = [shape._element for shape in shapes]
    for element in elements:
        tree.remove(element)
    if position == "front":
        for element in elements:
            tree.insert_element_before(element, "p:extLst")
    elif position == "back":
        for offset, element in enumerate(elements, start=2):
            tree.insert(offset, element)
    else:
        index = modifier.get("index")
        if isinstance(index, bool) or not isinstance(index, int) or index < 0:
            raise ValueError("reorder_elements requires position front/back or a non-negative index")
        for offset, element in enumerate(elements):
            tree.insert(index + 2 + offset, element)


def _align_elements(shapes, modifier: dict[str, Any]) -> None:
    alignment = modifier.get("alignment")
    bounds = (
        min(shape.left for shape in shapes), min(shape.top for shape in shapes),
        max(shape.left + shape.width for shape in shapes),
        max(shape.top + shape.height for shape in shapes),
    )
    for shape in shapes:
        if alignment == "left":
            shape.left = bounds[0]
        elif alignment == "center":
            shape.left = (bounds[0] + bounds[2] - shape.width) // 2
        elif alignment == "right":
            shape.left = bounds[2] - shape.width
        elif alignment == "top":
            shape.top = bounds[1]
        elif alignment == "middle":
            shape.top = (bounds[1] + bounds[3] - shape.height) // 2
        elif alignment == "bottom":
            shape.top = bounds[3] - shape.height
        else:
            raise ValueError("alignment must be left, center, right, top, middle, or bottom")


def _distribute_elements(shapes, modifier: dict[str, Any]) -> None:
    direction = modifier.get("direction")
    if direction == "horizontal":
        shapes.sort(key=lambda shape: shape.left)
        gap = (shapes[-1].left + shapes[-1].width - shapes[0].left - sum(shape.width for shape in shapes)) // (len(shapes) - 1)
        cursor = shapes[0].left + shapes[0].width
        for shape in shapes[1:-1]:
            shape.left = cursor + gap
            cursor = shape.left + shape.width
    elif direction == "vertical":
        shapes.sort(key=lambda shape: shape.top)
        gap = (shapes[-1].top + shapes[-1].height - shapes[0].top - sum(shape.height for shape in shapes)) // (len(shapes) - 1)
        cursor = shapes[0].top + shapes[0].height
        for shape in shapes[1:-1]:
            shape.top = cursor + gap
            cursor = shape.top + shape.height
    else:
        raise ValueError("direction must be horizontal or vertical")


def _slide_from_modifier(modifier: dict[str, Any]) -> Slide:
    slide = modifier.get("slide")
    if not isinstance(slide, Slide):
        raise TypeError("modifier.slide must be a Slide")
    return slide


def apply_modifiers(presentation: Presentation, modifiers: Sequence[dict[str, Any]]) -> None:
    """Apply ordered JSON-level modifiers to a loaded presentation."""
    for number, modifier in enumerate(modifiers, start=1):
        if not isinstance(modifier, dict):
            raise TypeError(f"modifier {number} must be an object")
        operation = modifier.get("operation")
        if operation == "append_slides":
            for spec in modifier.get("slides", ()):
                _add_slide(presentation, spec)
        elif operation == "insert_slides":
            index = modifier.get("index")
            if isinstance(index, bool) or not isinstance(index, int):
                raise TypeError("insert_slides.index must be an integer")
            for offset, spec in enumerate(modifier.get("slides", ())):
                _add_slide(presentation, spec, index + offset)
        elif operation == "replace_slide":
            index, slide = _selected_slide(presentation, modifier)
            spec = _slide_from_modifier(modifier)
            _clear_slide(slide)
            mapping = {}
            for element in spec.elements:
                shape = add_element(slide, element)
                mapping[element.id] = shape.shape_id
            apply_slide_effects(slide, spec, mapping, index + 1)
        elif operation == "delete_slides":
            has_indices = "indices" in modifier
            has_selectors = "selectors" in modifier
            if has_indices == has_selectors:
                raise ValueError("delete_slides requires exactly one of indices or selectors")
            if has_selectors:
                raw = modifier["selectors"]
                if not isinstance(raw, list):
                    raise TypeError("delete_slides.selectors must be an array")
                parsed = [
                    _selector(item, f"selectors[{position}]", shape=False)
                    for position, item in enumerate(raw)
                ]
                slide_ids = [item["slide_id"] for item in parsed]
                if len(slide_ids) != len(set(slide_ids)):
                    raise ValueError("delete_slides.selectors must be unique")
                indices = [_slide_by_id(presentation, slide_id)[0] for slide_id in slide_ids]
            else:
                indices = modifier["indices"]
                if not isinstance(indices, list) or not all(
                    isinstance(i, int) and not isinstance(i, bool) for i in indices
                ):
                    raise TypeError("delete_slides.indices must be an integer array")
            for index in sorted(set(indices), reverse=True):
                _slide_at(presentation, index)
                slide_id = presentation.slides._sldIdLst[index]
                presentation.part.drop_rel(slide_id.rId)
                presentation.slides._sldIdLst.remove(slide_id)
        elif operation == "move_slide":
            has_index = "from_index" in modifier
            has_selector = "from_selector" in modifier
            if has_index == has_selector:
                raise ValueError("move_slide requires exactly one of from_index or from_selector")
            if has_selector:
                parsed = _selector(modifier["from_selector"], "from_selector", shape=False)
                source, _slide = _slide_by_id(presentation, parsed["slide_id"])
            else:
                source = modifier["from_index"]
                _slide_at(presentation, source, "from_index")
            target = modifier.get("to_index")
            if isinstance(target, bool) or not isinstance(target, int) or not 0 <= target < len(presentation.slides):
                raise IndexError(f"to_index {target} is out of range")
            ids = presentation.slides._sldIdLst
            slide_id = ids[source]
            ids.remove(slide_id)
            ids.insert(target, slide_id)
        elif operation == "clear_slide":
            index, slide = _selected_slide(presentation, modifier)
            _clear_slide(slide)
            apply_slide_effects(slide, Slide(), {}, index + 1)
        elif operation == "add_elements":
            _index, slide = _selected_slide(presentation, modifier)
            for element in modifier.get("elements", ()):
                if element.id in _shape_names(slide):
                    raise ValueError(f"duplicate element ID on slide: {element.id}")
                add_element(slide, element)
        elif operation == "update_element":
            _index, _slide, shape = _selected_shape(presentation, modifier)
            _update_element(shape, modifier.get("changes"))
        elif operation == "duplicate_elements":
            _index, slide, shapes = _selected_shapes(presentation, modifier)
            _duplicate_elements(slide, shapes, modifier)
        elif operation == "reorder_elements":
            _index, slide, shapes = _selected_shapes(presentation, modifier)
            _reorder_elements(slide, shapes, modifier)
        elif operation == "align_elements":
            _index, _slide, shapes = _selected_shapes(presentation, modifier, minimum=2)
            _align_elements(shapes, modifier)
        elif operation == "distribute_elements":
            _index, _slide, shapes = _selected_shapes(presentation, modifier, minimum=3)
            _distribute_elements(shapes, modifier)
        elif operation == "replace_element":
            _index, slide, old = _selected_shape(presentation, modifier)
            _require_top_level(slide, old, "replace_element")
            old_shape_id = old.shape_id
            old_name = old.name
            element = modifier.get("element")
            if not isinstance(element, Element):
                raise TypeError("replace_element.element must be an Element")
            slide.shapes._spTree.remove(old._element)
            replacement = add_element(slide, element)
            replacements = (
                {old_shape_id: replacement.shape_id} if element.id == old_name else {}
            )
            update_timing_targets(
                slide,
                {old_shape_id},
                replacements=replacements,
                policy=modifier.get("animation_policy", "remove_targeted"),
            )
        elif operation == "delete_elements":
            _index, slide, shapes = _selected_shapes(presentation, modifier)
            for shape in shapes:
                _require_top_level(slide, shape, "delete_elements")
            removed_shape_ids = {shape.shape_id for shape in shapes}
            update_timing_targets(
                slide,
                removed_shape_ids,
                policy=modifier.get("animation_policy", "remove_targeted"),
            )
            for shape in shapes:
                slide.shapes._spTree.remove(shape._element)
        elif operation in ("set_animations", "set_transition"):
            index, slide = _selected_slide(presentation, modifier)
            animations = modifier.get("animations", ()) if operation == "set_animations" else ()
            transition = modifier.get("transition") if operation == "set_transition" else None
            spec = Slide(animations=tuple(animations), transition=transition)
            mapping = _mapping(slide) if operation == "set_animations" else {}
            apply_slide_effects(slide, spec, mapping, index + 1, replace_animations=operation == "set_animations", replace_transition=operation == "set_transition")
        else:
            raise ValueError(f"unsupported PowerPoint modifier operation: {operation!r}")


def edit_presentation(input_path: str | Path, output_path: str | Path, modifiers: Sequence[dict[str, Any]]) -> Path:
    """Edit an existing deck atomically, preserving untouched slides and content."""
    source, destination = Path(input_path), Path(output_path)
    if not source.is_file():
        raise FileNotFoundError(f"input presentation does not exist: {source}")
    if destination.suffix.lower() != ".pptx":
        raise ValueError("output_path must have a .pptx extension")
    presentation = Presentation(str(source))
    apply_modifiers(presentation, modifiers)
    destination.parent.mkdir(parents=True, exist_ok=True)
    temporary_path: Path | None = None
    try:
        with NamedTemporaryFile(prefix=f".{destination.stem}-", suffix=".pptx", dir=destination.parent, delete=False) as temporary:
            temporary_path = Path(temporary.name)
        presentation.save(str(temporary_path))
        validate_package(temporary_path, expected_slide_count=len(presentation.slides))
        os.replace(temporary_path, destination)
        temporary_path = None
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)
    return destination
