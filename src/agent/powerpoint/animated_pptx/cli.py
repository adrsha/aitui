"""Private JSON process bridge used only by AiTUI's native PowerPoint tool."""

from __future__ import annotations

import json
import sys
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any

from pptx import Presentation

from .editor import edit_presentation
from .generator import generate_presentation
from .inspect import inspect_presentation
from .model import Animation, Element, Slide
from .package_editor import PACKAGE_OPERATIONS, apply_package_modifiers


def _element(raw: Any, context: str) -> Element:
    if not isinstance(raw, dict):
        raise TypeError(f"{context} must be an object")
    return Element(**raw)


def _animation(raw: Any, context: str) -> Animation:
    if not isinstance(raw, dict):
        raise TypeError(f"{context} must be an object")
    return Animation(**raw)


def _slide(raw: Any, context: str) -> Slide:
    if not isinstance(raw, dict):
        raise TypeError(f"{context} must be an object")
    raw_elements = raw.get("elements", [])
    raw_animations = raw.get("animations", [])
    if not isinstance(raw_elements, list):
        raise TypeError(f"{context}.elements must be an array")
    if not isinstance(raw_animations, list):
        raise TypeError(f"{context}.animations must be an array")
    return Slide(
        elements=tuple(_element(item, f"{context}.elements") for item in raw_elements),
        animations=tuple(_animation(item, f"{context}.animations") for item in raw_animations),
        transition=raw.get("transition"),
    )


def _slides(raw: Any, context: str = "slides") -> tuple[Slide, ...]:
    if not isinstance(raw, list):
        raise TypeError(f"{context} must be an array")
    return tuple(_slide(item, f"{context}[{index}]") for index, item in enumerate(raw))


def presentation_from_spec(spec: dict[str, Any]) -> tuple[tuple[Slide, ...], Path]:
    """Convert a create request into validated model objects."""
    if not isinstance(spec, dict):
        raise TypeError("PowerPoint request must be a JSON object")
    output_path = spec.get("output_path")
    if not isinstance(output_path, str) or not output_path.strip():
        raise ValueError("output_path must be a non-empty string")
    return _slides(spec.get("slides")), Path(output_path)


def _modifier(raw: Any, index: int) -> dict[str, Any]:
    if not isinstance(raw, dict):
        raise TypeError(f"modifier {index} must be an object")
    parsed = dict(raw)
    operation = parsed.get("operation")
    if not isinstance(operation, str):
        raise TypeError(f"modifier {index}.operation must be a string")
    if operation in ("append_slides", "insert_slides"):
        parsed["slides"] = _slides(parsed.get("slides"), f"modifier {index}.slides")
    elif operation == "replace_slide":
        parsed["slide"] = _slide(parsed.get("slide"), f"modifier {index}.slide")
    elif operation == "add_elements":
        elements = parsed.get("elements")
        if not isinstance(elements, list):
            raise TypeError(f"modifier {index}.elements must be an array")
        parsed["elements"] = tuple(
            _element(item, f"modifier {index}.elements") for item in elements
        )
    elif operation == "replace_element":
        parsed["element"] = _element(
            parsed.get("element"), f"modifier {index}.element"
        )
    elif operation == "set_animations":
        animations = parsed.get("animations")
        if not isinstance(animations, list):
            raise TypeError(f"modifier {index}.animations must be an array")
        parsed["animations"] = tuple(
            _animation(item, f"modifier {index}.animations") for item in animations
        )
    return parsed


def _package_modifiers(spec: dict[str, Any]) -> tuple[dict[str, Any], ...]:
    raw = spec.get("package_modifiers", [])
    if not isinstance(raw, list):
        raise TypeError("package_modifiers must be an array")
    for index, modifier in enumerate(raw, start=1):
        if not isinstance(modifier, dict):
            raise TypeError(f"package modifier {index} must be an object")
        if modifier.get("operation") not in PACKAGE_OPERATIONS:
            raise ValueError(f"unsupported package modifier operation: {modifier.get('operation')!r}")
    return tuple(raw)


def execute_spec(spec: dict[str, Any]) -> tuple[Path, int, str]:
    """Execute a friendly deck operation, then optional advanced package edits."""
    if not isinstance(spec, dict):
        raise TypeError("PowerPoint request must be a JSON object")
    operation = spec.get("operation", "create")
    output = spec.get("output_path")
    if not isinstance(output, str) or not output.strip():
        raise ValueError("output_path must be a non-empty string")
    output_path = Path(output)
    package_modifiers = _package_modifiers(spec)
    working_output = output_path
    temporary_path: Path | None = None
    if package_modifiers:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with NamedTemporaryFile(
            prefix=f".{output_path.stem}-friendly-", suffix=".pptx",
            dir=output_path.parent, delete=False,
        ) as temporary:
            temporary_path = Path(temporary.name)
        working_output = temporary_path
    try:
        if operation in ("create", "replace"):
            slides = _slides(spec.get("slides"))
            generated = generate_presentation(slides, working_output)
            count = len(slides)
        elif operation == "append":
            input_path = Path(spec.get("input_path", output))
            slides = _slides(spec.get("slides"))
            generated = edit_presentation(
                input_path, working_output,
                ({"operation": "append_slides", "slides": slides},),
            )
            count = len(Presentation(generated).slides)
        elif operation == "edit":
            input_path = spec.get("input_path", output)
            if not isinstance(input_path, str) or not input_path.strip():
                raise ValueError("input_path must be a non-empty string")
            raw_modifiers = spec.get("modifiers", [])
            if not isinstance(raw_modifiers, list):
                raise TypeError("modifiers must be an array")
            if not raw_modifiers and not package_modifiers:
                raise ValueError("edit requires modifiers or package_modifiers")
            modifiers = tuple(_modifier(raw, index) for index, raw in enumerate(raw_modifiers, 1))
            if modifiers:
                generated = edit_presentation(input_path, working_output, modifiers)
                count = len(Presentation(generated).slides)
            else:
                generated = Path(input_path)
                count = len(Presentation(generated).slides)
        else:
            raise ValueError("operation must be create, replace, append, or edit")
        if package_modifiers:
            generated = apply_package_modifiers(
                generated, output_path, package_modifiers, expected_slide_count=count
            )
        return generated, count, operation
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)


def main() -> int:
    """Read one request from stdin and emit a JSON result."""
    try:
        spec = json.load(sys.stdin)
        if isinstance(spec, dict) and spec.get("operation") == "inspect":
            input_path = spec.get("input_path")
            if not isinstance(input_path, str) or not input_path.strip():
                raise ValueError("inspect requires a non-empty input_path")
            result = inspect_presentation(input_path)
            print(json.dumps({"ok": True, "inspection": result}))
            return 0
        generated, count, _operation = execute_spec(spec)
        print(json.dumps({"ok": True, "path": str(generated), "slides": count}))
        return 0
    except Exception as error:  # CLI boundary: return a concise tool-safe failure.
        print(json.dumps({"ok": False, "error": str(error)}), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
