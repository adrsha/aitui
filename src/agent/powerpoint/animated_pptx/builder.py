"""Translate the content model into python-pptx objects."""

from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from .model import Element, Slide

SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5

_SHAPE_TYPES = {
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
}


@dataclass(frozen=True, slots=True)
class BuiltPresentation:
    """Presentation plus element-name to shape-ID mappings for animation."""

    presentation: Presentation
    shape_ids: tuple[dict[str, int], ...]


def add_element(slide, element: Element):
    left, top = Inches(element.x), Inches(element.y)
    width, height = Inches(element.width), Inches(element.height)
    if element.type == "text":
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.text = element.text or ""
    elif element.type == "image":
        shape = slide.shapes.add_picture(
            str(element.image_path), left, top, width=width, height=height
        )
    else:
        shape = slide.shapes.add_shape(
            _SHAPE_TYPES[element.shape_type], left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(element.fill_color)
        shape.line.color.rgb = RGBColor.from_string(element.fill_color)
        if element.text is not None:
            shape.text = element.text

    shape.name = element.id
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(element.font_size)
                run.font.color.rgb = RGBColor.from_string(element.text_color)
    return shape


def build_presentation(slides: Sequence[Slide]) -> BuiltPresentation:
    """Build static slide content and retain exact shape IDs for animation."""
    if isinstance(slides, (str, bytes)) or not isinstance(slides, Sequence):
        raise TypeError("slides must be a sequence of Slide objects")
    if not all(isinstance(slide, Slide) for slide in slides):
        raise TypeError("slides must contain only Slide objects")

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_layout = presentation.slide_layouts[6]
    mappings: list[dict[str, int]] = []
    for slide_spec in slides:
        pptx_slide = presentation.slides.add_slide(blank_layout)
        mapping: dict[str, int] = {}
        for element in slide_spec.elements:
            shape = add_element(pptx_slide, element)
            mapping[element.id] = shape.shape_id
        mappings.append(mapping)
    return BuiltPresentation(presentation, tuple(mappings))
