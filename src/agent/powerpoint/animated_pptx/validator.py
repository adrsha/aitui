"""Post-save structural, OPC graph, and content validation for presentations."""

from __future__ import annotations

import posixpath
from pathlib import Path, PurePosixPath
from typing import Sequence
from urllib.parse import unquote, urlsplit
from zipfile import BadZipFile, ZipFile

from lxml import etree
from pptx import Presentation

from .model import Slide

PML_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"
RELATIONSHIPS_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/content-types"
NS = {"p": PML_NAMESPACE}
_XML_PARSER = etree.XMLParser(resolve_entities=False, no_network=True)


class PresentationValidationError(ValueError):
    """Raised when a generated presentation fails round-trip validation."""


def _parse_xml(data: bytes, member: str) -> etree._Element:
    try:
        return etree.fromstring(data, parser=_XML_PARSER)
    except etree.XMLSyntaxError as error:
        raise PresentationValidationError(
            f"package XML member is malformed: {member}: {error}"
        ) from error


def _relationship_source(rels_member: str) -> str:
    if rels_member == "_rels/.rels":
        return ""
    path = PurePosixPath(rels_member)
    if path.parent.name != "_rels" or not path.name.endswith(".rels"):
        raise PresentationValidationError(f"invalid relationship part name: {rels_member}")
    source_name = path.name[:-5]
    return str(path.parent.parent / source_name)


def _internal_target(source: str, target: str, rels_member: str) -> str:
    parsed = urlsplit(target)
    if parsed.scheme or parsed.netloc:
        raise PresentationValidationError(
            f"internal relationship uses an external URI in {rels_member}: {target}"
        )
    decoded = unquote(parsed.path)
    if not decoded or "\\" in decoded or "\x00" in decoded:
        raise PresentationValidationError(
            f"invalid internal relationship target in {rels_member}: {target}"
        )
    if decoded.startswith("/"):
        normalized = posixpath.normpath(decoded).lstrip("/")
    else:
        normalized = posixpath.normpath(posixpath.join(posixpath.dirname(source), decoded))
    if normalized in ("", ".", "..") or normalized.startswith("../"):
        raise PresentationValidationError(
            f"relationship target escapes the package in {rels_member}: {target}"
        )
    return normalized


def _validate_content_types(members: set[str], root: etree._Element) -> None:
    expected_root = f"{{{CONTENT_TYPES_NAMESPACE}}}Types"
    if root.tag != expected_root:
        raise PresentationValidationError("[Content_Types].xml has an invalid root element")
    defaults: dict[str, str] = {}
    overrides: dict[str, str] = {}
    for child in root:
        if child.tag == f"{{{CONTENT_TYPES_NAMESPACE}}}Default":
            extension, content_type = child.get("Extension"), child.get("ContentType")
            if not extension or not content_type:
                raise PresentationValidationError("invalid Default in [Content_Types].xml")
            key = extension.lower()
            if key in defaults:
                raise PresentationValidationError(f"duplicate content type default: {extension}")
            defaults[key] = content_type
        elif child.tag == f"{{{CONTENT_TYPES_NAMESPACE}}}Override":
            part, content_type = child.get("PartName"), child.get("ContentType")
            if not part or not part.startswith("/") or not content_type:
                raise PresentationValidationError("invalid Override in [Content_Types].xml")
            key = part.lstrip("/")
            if key in overrides:
                raise PresentationValidationError(f"duplicate content type override: {part}")
            if key not in members:
                raise PresentationValidationError(
                    f"content type override targets a missing part: {part}"
                )
            overrides[key] = content_type
        else:
            raise PresentationValidationError("unexpected element in [Content_Types].xml")
    for member in members:
        if member == "[Content_Types].xml" or member.endswith("/"):
            continue
        name = PurePosixPath(member).name
        extension = "rels" if name.endswith(".rels") else PurePosixPath(member).suffix.lower().lstrip(".")
        if member not in overrides and extension not in defaults:
            raise PresentationValidationError(
                f"package part has no content type declaration: {member}"
            )


def _validate_relationships(
    members: set[str], xml_roots: dict[str, etree._Element]
) -> None:
    for rels_member, root in xml_roots.items():
        if not rels_member.endswith(".rels"):
            continue
        if root.tag != f"{{{RELATIONSHIPS_NAMESPACE}}}Relationships":
            raise PresentationValidationError(
                f"relationship part has an invalid root element: {rels_member}"
            )
        source = _relationship_source(rels_member)
        if source and source not in members:
            raise PresentationValidationError(
                f"relationship part has no source package part: {rels_member}"
            )
        ids: set[str] = set()
        for relationship in root:
            if relationship.tag != f"{{{RELATIONSHIPS_NAMESPACE}}}Relationship":
                raise PresentationValidationError(
                    f"unexpected element in relationship part: {rels_member}"
                )
            relationship_id = relationship.get("Id")
            relationship_type = relationship.get("Type")
            target = relationship.get("Target")
            if not relationship_id or not relationship_type or not target:
                raise PresentationValidationError(
                    f"relationship is missing Id, Type, or Target: {rels_member}"
                )
            if relationship_id in ids:
                raise PresentationValidationError(
                    f"duplicate relationship ID in {rels_member}: {relationship_id}"
                )
            ids.add(relationship_id)
            target_mode = relationship.get("TargetMode")
            if target_mode not in (None, "Internal", "External"):
                raise PresentationValidationError(
                    f"invalid TargetMode in {rels_member}: {target_mode}"
                )
            if target_mode != "External":
                resolved = _internal_target(source, target, rels_member)
                if resolved not in members:
                    raise PresentationValidationError(
                        f"dangling relationship in {rels_member}: {relationship_id} -> {resolved}"
                    )


def validate_package(path: str | Path, expected_slide_count: int | None = None) -> None:
    """Validate reopenability, every XML member, content types, and relationships."""
    file_path = Path(path)
    if not file_path.is_file():
        raise FileNotFoundError(f"presentation does not exist: {file_path}")
    try:
        with ZipFile(file_path) as archive:
            infos = archive.infolist()
            names = [info.filename for info in infos]
            if len(names) != len(set(names)):
                raise PresentationValidationError("presentation contains duplicate ZIP members")
            bad_member = archive.testzip()
            if bad_member is not None:
                raise PresentationValidationError(
                    f"corrupt ZIP member in presentation: {bad_member}"
                )
            members = set(names)
            if "[Content_Types].xml" not in members:
                raise PresentationValidationError("presentation is missing [Content_Types].xml")
            xml_roots = {
                name: _parse_xml(archive.read(name), name)
                for name in names
                if name.endswith(".xml") or name.endswith(".rels")
            }
            _validate_content_types(members, xml_roots["[Content_Types].xml"])
            _validate_relationships(members, xml_roots)
    except BadZipFile as error:
        raise PresentationValidationError(f"invalid pptx ZIP package: {error}") from error

    try:
        reopened = Presentation(str(file_path))
    except Exception as error:
        raise PresentationValidationError(
            f"python-pptx could not reopen {file_path}: {error}"
        ) from error
    if expected_slide_count is not None and len(reopened.slides) != expected_slide_count:
        raise PresentationValidationError(
            f"slide count mismatch: expected {expected_slide_count}, got {len(reopened.slides)}"
        )


def validate_presentation(path: str | Path, slides: Sequence[Slide]) -> None:
    """Reopen generated content and verify its modeled slide-level semantics."""
    if not isinstance(path, (str, Path)):
        raise TypeError("path must be a string or pathlib.Path")
    if isinstance(slides, (str, bytes)) or not isinstance(slides, Sequence):
        raise TypeError("slides must be a sequence of Slide objects")
    if not all(isinstance(slide, Slide) for slide in slides):
        raise TypeError("slides must contain only Slide objects")
    file_path = Path(path)
    validate_package(file_path, expected_slide_count=len(slides))
    reopened = Presentation(str(file_path))
    for index, (actual, expected) in enumerate(zip(reopened.slides, slides), start=1):
        if len(actual.shapes) != len(expected.elements):
            raise PresentationValidationError(
                f"slide {index} shape count mismatch: expected "
                f"{len(expected.elements)}, got {len(actual.shapes)}"
            )
        actual_text = [shape.text for shape in actual.shapes if shape.has_text_frame]
        expected_text = [
            element.text or "" for element in expected.elements
            if element.type in ("text", "shape")
        ]
        if actual_text != expected_text:
            raise PresentationValidationError(
                f"slide {index} text mismatch: expected {expected_text!r}, got {actual_text!r}"
            )

    with ZipFile(file_path) as archive:
        for index, slide in enumerate(slides, start=1):
            member = f"ppt/slides/slide{index}.xml"
            try:
                root = _parse_xml(archive.read(member), member)
            except KeyError as error:
                raise PresentationValidationError(
                    f"slide {index} XML is missing: {error}"
                ) from error
            if slide.animations and root.find("p:timing", NS) is None:
                raise PresentationValidationError(
                    f"slide {index} is missing required p:timing animation XML"
                )
            if slide.transition and root.find("p:transition", NS) is None:
                raise PresentationValidationError(
                    f"slide {index} is missing required p:transition XML"
                )
