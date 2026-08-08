"""Inject a fixed, deliberately small set of PowerPoint timing templates."""

from __future__ import annotations

from collections.abc import Sequence

from lxml import etree
from pptx import Presentation
from pptx.presentation import Presentation as PresentationObject

from .model import Animation, Slide

PML_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"
P14_NAMESPACE = "http://schemas.microsoft.com/office/powerpoint/2010/main"
NS = {"p": PML_NAMESPACE}

# These are fixed ECMA-376 timing behaviors, not an extensible animation DSL.
# preset IDs match PowerPoint's built-in effect catalog; the filter selects the
# validated animEffect behavior and direction.
_EFFECTS: dict[str, tuple[str, str, str, str]] = {
    "fade_in": ("entr", "10", "0", "fade"),
    "fly_in_left": ("entr", "2", "8", "slide(fromLeft)"),
    "fly_in_right": ("entr", "2", "2", "slide(fromRight)"),
    "fly_in_bottom": ("entr", "2", "4", "slide(fromBottom)"),
    "wipe": ("entr", "22", "1", "wipe(right)"),
    "zoom": ("entr", "23", "0", "zoom(in)"),
    "fade_out": ("exit", "10", "0", "fade"),
}

_TRANSITIONS: dict[str, tuple[str, dict[str, str]]] = {
    "fade": ("fade", {}),
    "push_left": ("push", {"dir": "l"}),
    "wipe_left": ("wipe", {"dir": "l"}),
}


class AnimationTargetError(ValueError):
    """Raised when an animation names no shape on its slide."""


def _p(tag: str) -> etree._Element:
    return etree.Element(f"{{{PML_NAMESPACE}}}{tag}")


def _sub(parent: etree._Element, tag: str, **attributes: str) -> etree._Element:
    return etree.SubElement(parent, f"{{{PML_NAMESPACE}}}{tag}", **attributes)


def _condition_list(parent: etree._Element, delay_ms: int) -> None:
    conditions = _sub(parent, "stCondLst")
    _sub(conditions, "cond", delay=str(delay_ms))


def _target(parent: etree._Element, shape_id: int) -> None:
    target = _sub(parent, "tgtEl")
    _sub(target, "spTgt", spid=str(shape_id))


def _effect_node(
    animation: Animation, shape_id: int, node_id: int
) -> tuple[etree._Element, int]:
    preset_class, preset_id, preset_subtype, effect_filter = _EFFECTS[animation.type]
    wrapper = _p("par")
    wrapper_tn = _sub(
        wrapper, "cTn", id=str(node_id), fill="hold", nodeType="clickEffect",
        presetClass=preset_class, presetID=preset_id, presetSubtype=preset_subtype,
    )
    node_id += 1
    _condition_list(wrapper_tn, animation.delay_ms)
    children = _sub(wrapper_tn, "childTnLst")
    effect = _sub(
        children, "animEffect",
        transition="out" if animation.type == "fade_out" else "in",
        filter=effect_filter,
    )
    behavior = _sub(effect, "cBhvr")
    _sub(behavior, "cTn", id=str(node_id), dur=str(animation.duration_ms), fill="hold")
    node_id += 1
    _target(behavior, shape_id)
    return wrapper, node_id


def _timing_node(animations: Sequence[tuple[Animation, int]]) -> etree._Element:
    timing = _p("timing")
    timing_list = _sub(timing, "tnLst")
    root_parallel = _sub(timing_list, "par")
    root_tn = _sub(
        root_parallel, "cTn", id="1", dur="indefinite", restart="never",
        nodeType="tmRoot",
    )
    root_children = _sub(root_tn, "childTnLst")
    main_sequence = _sub(root_children, "seq", concurrent="1", nextAc="seek")
    main_tn = _sub(
        main_sequence, "cTn", id="2", dur="indefinite", nodeType="mainSeq"
    )
    main_children = _sub(main_tn, "childTnLst")

    node_id = 3
    for animation, shape_id in animations:
        effect, node_id = _effect_node(animation, shape_id, node_id)
        # PowerPoint uses clickEffect for all sequence entries. Trigger grouping
        # is represented by the start condition: indefinite means user click;
        # finite delay means automatic relative to the preceding entry.
        effect_tn = effect.find("p:cTn", NS)
        if effect_tn is None:  # internal invariant
            raise RuntimeError("fixed animation template has no cTn node")
        condition = effect_tn.find("p:stCondLst/p:cond", NS)
        if condition is None:  # internal invariant
            raise RuntimeError("fixed animation template has no start condition")
        if animation.trigger == "on_click":
            condition.set("delay", "indefinite")
        elif animation.trigger == "with_previous":
            effect_tn.set("nodeType", "withEffect")
        else:
            effect_tn.set("nodeType", "afterEffect")
        main_children.append(effect)

    previous = _sub(main_sequence, "prevCondLst")
    prev_condition = _sub(previous, "cond", evt="onPrev", delay="0")
    prev_target = _sub(prev_condition, "tgtEl")
    _sub(prev_target, "sldTgt")
    following = _sub(main_sequence, "nextCondLst")
    next_condition = _sub(following, "cond", evt="onNext", delay="0")
    next_target = _sub(next_condition, "tgtEl")
    _sub(next_target, "sldTgt")
    return timing


def _insert_before_extension_list(slide_root: etree._Element, node: etree._Element) -> None:
    extension = slide_root.find("p:extLst", NS)
    if extension is None:
        slide_root.append(node)
    else:
        extension.addprevious(node)


def apply_animations(
    presentation: Presentation,
    slides: Sequence[Slide],
    shape_ids: Sequence[dict[str, int]],
) -> None:
    """Inject fixed timing and transition XML into a built presentation."""
    if not isinstance(presentation, PresentationObject):
        raise TypeError("presentation must be a python-pptx Presentation")
    if isinstance(slides, (str, bytes)) or not isinstance(slides, Sequence):
        raise TypeError("slides must be a sequence of Slide objects")
    if not all(isinstance(slide, Slide) for slide in slides):
        raise TypeError("slides must contain only Slide objects")
    if isinstance(shape_ids, (str, bytes)) or not isinstance(shape_ids, Sequence):
        raise TypeError("shape_ids must be a sequence of dictionaries")
    if len(presentation.slides) != len(slides) or len(shape_ids) != len(slides):
        raise ValueError("presentation, slides, and shape_ids must have equal lengths")

    for index, (pptx_slide, slide_spec, mapping) in enumerate(
        zip(presentation.slides, slides, shape_ids), start=1
    ):
        if not isinstance(mapping, dict) or not all(
            isinstance(name, str) and isinstance(shape_id, int)
            for name, shape_id in mapping.items()
        ):
            raise TypeError("each shape_ids item must map strings to integer shape IDs")
        actual_ids = {shape.shape_id for shape in pptx_slide.shapes}
        resolved: list[tuple[Animation, int]] = []
        for animation in sorted(slide_spec.animations, key=lambda item: item.order):
            shape_id = mapping.get(animation.target)
            if shape_id is None:
                raise AnimationTargetError(
                    f"slide {index} animation target {animation.target!r} does not exist"
                )
            if shape_id not in actual_ids:
                raise AnimationTargetError(
                    f"slide {index} animation target {animation.target!r} maps to missing "
                    f"shape ID {shape_id}"
                )
            resolved.append((animation, shape_id))

        root = pptx_slide._element
        existing_timing = root.find("p:timing", NS)
        if existing_timing is not None:
            root.remove(existing_timing)
        existing_transition = root.find("p:transition", NS)
        if existing_transition is not None:
            root.remove(existing_transition)
        if slide_spec.transition is not None:
            child_tag, child_attributes = _TRANSITIONS[slide_spec.transition]
            transition = _p("transition")
            transition.set("spd", "med")
            _sub(transition, child_tag, **child_attributes)
            _insert_before_extension_list(root, transition)
        if resolved:
            _insert_before_extension_list(root, _timing_node(resolved))
