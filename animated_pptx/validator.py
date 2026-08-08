"""Post-save structural and content validation for generated presentations."""

from __future__ import annotations

from pathlib import Path
from typing import Sequence
from zipfile import BadZipFile, ZipFile

from lxml import etree
from pptx import Presentation

from .model import Slide

PML_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"p": PML_NAMESPACE}


class PresentationValidationError(ValueError):
    """Raised when a generated presentation fails round-trip validation."""


def validate_presentation(path: str | Path, slides: Sequence[Slide]) -> None:
    """Re-open with python-pptx and parse every slide XML with lxml."""
    if not isinstance(path, (str, Path)):
        raise TypeError("path must be a string or pathlib.Path")
    if isinstance(slides, (str, bytes)) or not isinstance(slides, Sequence):
        raise TypeError("slides must be a sequence of Slide objects")
    if not all(isinstance(slide, Slide) for slide in slides):
        raise TypeError("slides must contain only Slide objects")
    file_path = Path(path)
    if not file_path.is_file():
        raise FileNotFoundError(f"presentation does not exist: {file_path}")

    try:
        reopened = Presentation(str(file_path))
    except (OSError, ValueError, KeyError) as error:
        raise PresentationValidationError(
            f"python-pptx could not reopen {file_path}: {error}"
        ) from error
    if len(reopened.slides) != len(slides):
        raise PresentationValidationError(
            f"slide count mismatch: expected {len(slides)}, got {len(reopened.slides)}"
        )
    for index, (actual, expected) in enumerate(zip(reopened.slides, slides), start=1):
        if len(actual.shapes) != len(expected.elements):
            raise PresentationValidationError(
                f"slide {index} shape count mismatch: expected "
                f"{len(expected.elements)}, got {len(actual.shapes)}"
            )
        actual_text = [shape.text for shape in actual.shapes if shape.has_text_frame]
        expected_text = [
            element.text or "" for element in expected.elements
            if element.type == "text" or (element.type == "shape" and element.text is not None)
        ]
        if actual_text != expected_text:
            raise PresentationValidationError(
                f"slide {index} text mismatch: expected {expected_text!r}, got {actual_text!r}"
            )

    try:
        with ZipFile(file_path) as archive:
            bad_member = archive.testzip()
            if bad_member is not None:
                raise PresentationValidationError(
                    f"corrupt ZIP member in presentation: {bad_member}"
                )
            for index, slide in enumerate(slides, start=1):
                member = f"ppt/slides/slide{index}.xml"
                try:
                    root = etree.fromstring(archive.read(member))
                except (KeyError, etree.XMLSyntaxError) as error:
                    raise PresentationValidationError(
                        f"slide {index} XML is missing or malformed: {error}"
                    ) from error
                if slide.animations and root.find("p:timing", NS) is None:
                    raise PresentationValidationError(
                        f"slide {index} is missing required p:timing animation XML"
                    )
                if slide.transition and root.find("p:transition", NS) is None:
                    raise PresentationValidationError(
                        f"slide {index} is missing required p:transition XML"
                    )
    except BadZipFile as error:
        raise PresentationValidationError(f"invalid pptx ZIP package: {error}") from error
