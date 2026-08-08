import json
import os
from pathlib import Path
import subprocess
import sys
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation
from pptx.util import Inches
import pytest

from animated_pptx import Animation, Element, Slide, generate_presentation, inspect_presentation
from animated_pptx.animator import AnimationTargetError, NS
from animated_pptx.cli import execute_spec, presentation_from_spec


def slide_xml(path: Path, number: int = 1):
    with ZipFile(path) as archive:
        return etree.fromstring(archive.read(f"ppt/slides/slide{number}.xml"))


def test_empty_slide(tmp_path: Path) -> None:
    output = generate_presentation((Slide(),), tmp_path / "empty-slide.pptx")
    reopened = Presentation(output)
    assert len(reopened.slides) == 1
    assert len(reopened.slides[0].shapes) == 0
    assert slide_xml(output).find("p:timing", NS) is None


def test_zero_slides(tmp_path: Path) -> None:
    output = generate_presentation((), tmp_path / "zero-slides.pptx")
    assert len(Presentation(output).slides) == 0


def test_missing_animation_target_raises_without_replacing_destination(tmp_path: Path) -> None:
    destination = tmp_path / "missing.pptx"
    destination.write_bytes(b"existing")
    slides = (
        Slide(
            elements=(Element("present", "text", 1, 1, 4, 1, text="Present"),),
            animations=(Animation("fade_in", "absent", 0),),
        ),
    )
    with pytest.raises(AnimationTargetError, match="slide 1.*'absent'.*does not exist"):
        generate_presentation(slides, destination)
    assert destination.read_bytes() == b"existing"


def test_full_multi_slide_mixed_animation_deck(tmp_path: Path) -> None:
    slides = (
        Slide(
            elements=(
                Element("a", "text", 0.5, 0.5, 4, 1, text="Alpha"),
                Element("b", "shape", 1, 2, 3, 2, text="Beta"),
                Element("c", "text", 5, 2, 4, 1, text="Gamma"),
            ),
            animations=(
                Animation("fade_in", "a", 0, duration_ms=400),
                Animation("fly_in_left", "b", 1, trigger="after_previous"),
                Animation("wipe", "c", 2, trigger="with_previous", delay_ms=100),
            ),
            transition="fade",
        ),
        Slide(
            elements=(
                Element("d", "shape", 1, 1, 3, 2, text="Delta"),
                Element("e", "shape", 5, 1, 3, 2, text="Epsilon"),
                Element("f", "text", 3, 4, 5, 1, text="Zeta"),
            ),
            animations=(
                Animation("fly_in_right", "d", 0),
                Animation("fly_in_bottom", "e", 1, trigger="after_previous"),
                Animation("zoom", "f", 2, trigger="after_previous"),
                Animation("fade_out", "f", 3, trigger="after_previous", delay_ms=250),
            ),
            transition="wipe_left",
        ),
    )
    output = generate_presentation(slides, tmp_path / "mixed.pptx")
    reopened = Presentation(output)
    assert len(reopened.slides) == 2
    assert [[shape.text for shape in slide.shapes] for slide in reopened.slides] == [
        ["Alpha", "Beta", "Gamma"], ["Delta", "Epsilon", "Zeta"]
    ]
    for number, expected_effects in ((1, 3), (2, 4)):
        root = slide_xml(output, number)
        assert root.find("p:timing", NS) is not None
        assert root.find("p:transition", NS) is not None
        effects = root.findall(".//p:animEffect", NS)
        assert len(effects) == expected_effects
        actual_shape_ids = {shape.shape_id for shape in reopened.slides[number - 1].shapes}
        referenced_ids = {
            int(target.get("spid")) for target in root.findall(".//p:spTgt", NS)
        }
        assert referenced_ids <= actual_shape_ids


def test_json_bridge_builds_typed_slide_model(tmp_path: Path) -> None:
    slides, output = presentation_from_spec({
        "output_path": str(tmp_path / "bridge.pptx"),
        "slides": [{
            "elements": [{
                "id": "title", "type": "text", "x": 1, "y": 1,
                "width": 5, "height": 1, "text": "Bridge",
            }],
            "animations": [{"type": "fade_in", "target": "title", "order": 0}],
            "transition": "fade",
        }],
    })
    assert output == tmp_path / "bridge.pptx"
    assert slides[0].elements[0].text == "Bridge"
    assert slides[0].animations[0].type == "fade_in"



def test_append_and_edit_existing_deck_with_simple_json_modifiers(tmp_path: Path) -> None:
    deck = tmp_path / "editable.pptx"
    execute_spec({
        "operation": "create",
        "output_path": str(deck),
        "slides": [{
            "elements": [{
                "id": "title", "type": "text", "x": 1, "y": 1,
                "width": 5, "height": 1, "text": "Original",
            }],
            "animations": [{"type": "fade_in", "target": "title", "order": 0}],
            "transition": "fade",
        }],
    })
    execute_spec({
        "operation": "append",
        "output_path": str(deck),
        "slides": [{
            "elements": [{
                "id": "second", "type": "shape", "x": 2, "y": 2,
                "width": 3, "height": 2, "text": "Second",
            }]
        }],
    })
    execute_spec({
        "operation": "edit",
        "output_path": str(deck),
        "modifiers": [
            {
                "operation": "update_element", "slide_index": 0,
                "element_id": "title",
                "changes": {"text": "Updated", "x": 1.5, "font_size": 30},
            },
            {
                "operation": "add_elements", "slide_index": 0,
                "elements": [{
                    "id": "badge", "type": "shape", "shape_type": "ellipse",
                    "x": 8, "y": 1,
                    "width": 1, "height": 1,
                }],
            },
            {"operation": "set_transition", "slide_index": 1, "transition": "wipe_left"},
            {"operation": "move_slide", "from_index": 1, "to_index": 0},
        ],
    })
    reopened = Presentation(deck)
    assert len(reopened.slides) == 2
    assert reopened.slides[0].shapes[0].text == "Second"
    assert [shape.name for shape in reopened.slides[1].shapes] == ["title", "badge"]
    assert reopened.slides[1].shapes[0].text == "Updated"
    with ZipFile(deck) as archive:
        moved_part = str(reopened.slides[0].part.partname).lstrip("/")
        moved_root = etree.fromstring(archive.read(moved_part))
    assert moved_root.find("p:transition", NS) is not None


def test_edit_can_insert_replace_delete_and_reset_animations(tmp_path: Path) -> None:
    deck = tmp_path / "crud.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [
            {"elements": [{"id": "a", "type": "text", "x": 1, "y": 1, "width": 3, "height": 1, "text": "A"}]},
            {"elements": [{"id": "b", "type": "text", "x": 1, "y": 1, "width": 3, "height": 1, "text": "B"}]},
        ],
    })
    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "modifiers": [
            {"operation": "insert_slides", "index": 1, "slides": [{"elements": [{"id": "i", "type": "text", "x": 1, "y": 1, "width": 3, "height": 1, "text": "Inserted"}]}]},
            {"operation": "replace_slide", "slide_index": 2, "slide": {"elements": [{"id": "r", "type": "text", "x": 1, "y": 1, "width": 3, "height": 1, "text": "Replacement"}]}},
            {"operation": "replace_element", "slide_index": 0, "element_id": "a", "element": {"id": "new-a", "type": "shape", "x": 1, "y": 1, "width": 3, "height": 1, "text": "New A"}},
            {"operation": "set_animations", "slide_index": 0, "animations": [{"type": "zoom", "target": "new-a", "order": 0}]},
            {"operation": "delete_slides", "indices": [1]},
        ],
    })
    reopened = Presentation(deck)
    assert [[shape.text for shape in slide.shapes] for slide in reopened.slides] == [["New A"], ["Replacement"]]
    assert slide_xml(deck, 1).find("p:timing", NS) is not None



def animation_target_ids(path: Path, number: int = 1) -> list[int]:
    return [
        int(target.get("spid"))
        for target in slide_xml(path, number).findall(".//p:spTgt", NS)
    ]


def test_element_deletion_preserves_unrelated_animations(tmp_path: Path) -> None:
    deck = tmp_path / "animation-delete.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{
            "elements": [
                {"id": "a", "type": "shape", "x": 1, "y": 1, "width": 2, "height": 1},
                {"id": "b", "type": "shape", "x": 4, "y": 1, "width": 2, "height": 1},
                {"id": "c", "type": "shape", "x": 7, "y": 1, "width": 2, "height": 1},
            ],
            "animations": [
                {"type": "fade_in", "target": "a", "order": 0},
                {"type": "zoom", "target": "b", "order": 1},
            ],
        }],
    })
    original = Presentation(deck)
    original_ids = {shape.name: shape.shape_id for shape in original.slides[0].shapes}

    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "modifiers": [{
            "operation": "delete_elements", "slide_index": 0, "element_ids": ["c"],
        }],
    })
    assert animation_target_ids(deck) == [original_ids["a"], original_ids["b"]]

    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "modifiers": [{
            "operation": "delete_elements", "slide_index": 0, "element_ids": ["a"],
            "animation_policy": "remove_targeted",
        }],
    })
    assert animation_target_ids(deck) == [original_ids["b"]]


def test_animation_reference_policy_fails_atomically(tmp_path: Path) -> None:
    deck = tmp_path / "animation-policy.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{
            "elements": [{"id": "a", "type": "shape", "x": 1, "y": 1, "width": 2, "height": 1}],
            "animations": [{"type": "fade_in", "target": "a", "order": 0}],
        }],
    })
    before = deck.read_bytes()
    with pytest.raises(AnimationTargetError, match="would remove animation target"):
        execute_spec({
            "operation": "edit", "output_path": str(deck),
            "modifiers": [{
                "operation": "delete_elements", "slide_index": 0,
                "element_ids": ["a"], "animation_policy": "error_if_referenced",
            }],
        })
    assert deck.read_bytes() == before


def test_same_id_replacement_retargets_animation_and_preserves_others(tmp_path: Path) -> None:
    deck = tmp_path / "animation-replace.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{
            "elements": [
                {"id": "a", "type": "shape", "x": 1, "y": 1, "width": 2, "height": 1},
                {"id": "b", "type": "shape", "x": 4, "y": 1, "width": 2, "height": 1},
            ],
            "animations": [
                {"type": "fade_in", "target": "a", "order": 0},
                {"type": "zoom", "target": "b", "order": 1},
            ],
        }],
    })
    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "modifiers": [{
            "operation": "replace_element", "slide_index": 0, "element_id": "a",
            "element": {"id": "a", "type": "shape", "shape_type": "ellipse", "x": 1, "y": 1, "width": 2, "height": 1},
        }],
    })
    reopened = Presentation(deck)
    ids = {shape.name: shape.shape_id for shape in reopened.slides[0].shapes}
    assert animation_target_ids(deck) == [ids["a"], ids["b"]]


def test_high_level_geometry_duplicate_and_z_order_modifiers(tmp_path: Path) -> None:
    deck = tmp_path / "geometry.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{"elements": [
            {"id": "a", "type": "shape", "x": 1, "y": 1, "width": 1, "height": 1},
            {"id": "b", "type": "shape", "x": 3, "y": 2, "width": 1, "height": 1},
            {"id": "c", "type": "shape", "x": 6, "y": 3, "width": 1, "height": 1},
        ]}],
    })
    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "modifiers": [
            {"operation": "align_elements", "slide_index": 0, "element_ids": ["a", "b", "c"], "alignment": "top"},
            {"operation": "distribute_elements", "slide_index": 0, "element_ids": ["a", "b", "c"], "direction": "horizontal"},
            {"operation": "duplicate_elements", "slide_index": 0, "element_ids": ["b"], "new_ids": ["b-copy"], "offset_x": 0.5, "offset_y": 0.25},
            {"operation": "update_element", "slide_index": 0, "element_id": "b-copy", "changes": {"rotation": 30, "flip_horizontal": True}},
            {"operation": "reorder_elements", "slide_index": 0, "element_ids": ["a"], "position": "front"},
        ],
    })
    shapes = {shape.name: shape for shape in Presentation(deck).slides[0].shapes}
    assert set(shapes) == {"a", "b", "c", "b-copy"}
    assert shapes["a"].top == shapes["b"].top == shapes["c"].top
    assert shapes["b-copy"].rotation == 30
    assert list(shapes)[-1] == "a"


def test_package_modifiers_patch_any_xml_part_atomically(tmp_path: Path) -> None:
    deck = tmp_path / "advanced.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{"elements": [{
            "id": "title", "type": "text", "x": 1, "y": 1,
            "width": 5, "height": 1, "text": "Before",
        }]}],
        "package_modifiers": [{
            "operation": "patch_xml", "part": "ppt/slides/slide1.xml",
            "xpath": ".//a:t", "namespaces": {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
            "action": "set_text", "text": "After",
        }],
    })
    assert Presentation(deck).slides[0].shapes[0].text == "After"

    before = deck.read_bytes()
    with pytest.raises(ValueError, match="matched no nodes"):
        execute_spec({
            "operation": "edit", "output_path": str(deck),
            "package_modifiers": [{
                "operation": "patch_xml", "part": "ppt/slides/slide1.xml",
                "xpath": ".//a:does-not-exist", "namespaces": {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
                "action": "remove",
            }],
        })
    assert deck.read_bytes() == before


def test_package_modifiers_manage_parts_relationships_and_content_types(tmp_path: Path) -> None:
    deck = tmp_path / "package-graph.pptx"
    execute_spec({"output_path": str(deck), "slides": [{}]})
    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "package_modifiers": [
            {
                "operation": "put_part", "part": "customXml/item1.xml",
                "xml": "<metadata xmlns=\"urn:aitui:test\"><value>kept</value></metadata>",
            },
            {
                "operation": "set_content_type", "part": "customXml/item1.xml",
                "content_type": "application/xml",
            },
            {
                "operation": "put_relationship", "source_part": "/", "id": "rIdAituiCustom",
                "relationship_type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/customXml",
                "target": "customXml/item1.xml",
            },
        ],
    })
    with ZipFile(deck) as archive:
        assert etree.fromstring(archive.read("customXml/item1.xml")).find(
            "{urn:aitui:test}value"
        ).text == "kept"
        relationships = etree.fromstring(archive.read("_rels/.rels"))
        assert relationships.xpath(
            "count(./pr:Relationship[@Id='rIdAituiCustom'])",
            namespaces={"pr": "http://schemas.openxmlformats.org/package/2006/relationships"},
        ) == 1
        content_types = etree.fromstring(archive.read("[Content_Types].xml"))
        assert content_types.xpath(
            "count(./ct:Override[@PartName='/customXml/item1.xml'])",
            namespaces={"ct": "http://schemas.openxmlformats.org/package/2006/content-types"},
        ) == 1


def test_package_validation_rejects_dangling_relationship_atomically(tmp_path: Path) -> None:
    deck = tmp_path / "dangling.pptx"
    execute_spec({"output_path": str(deck), "slides": [{}]})
    before = deck.read_bytes()
    with pytest.raises(ValueError, match="dangling relationship"):
        execute_spec({
            "operation": "edit", "output_path": str(deck),
            "package_modifiers": [{
                "operation": "put_relationship", "source_part": "/", "id": "rIdBroken",
                "relationship_type": "urn:aitui:missing", "target": "missing/item.xml",
            }],
        })
    assert deck.read_bytes() == before


def test_package_validation_checks_non_slide_xml_and_safe_paths(tmp_path: Path) -> None:
    deck = tmp_path / "invalid-package.pptx"
    execute_spec({"output_path": str(deck), "slides": [{}]})
    before = deck.read_bytes()
    with pytest.raises(ValueError, match="must stay inside"):
        execute_spec({
            "operation": "edit", "output_path": str(deck),
            "package_modifiers": [{
                "operation": "put_part", "part": "../outside.xml", "xml": "<safe/>",
            }],
        })
    assert deck.read_bytes() == before

    with pytest.raises(ValueError, match="malformed"):
        execute_spec({
            "operation": "edit", "output_path": str(deck),
            "package_modifiers": [
                {
                    "operation": "put_part", "part": "customXml/item1.xml",
                    "text": "<broken>",
                },
                {
                    "operation": "set_content_type", "part": "customXml/item1.xml",
                    "content_type": "application/xml",
                },
            ],
        })
    assert deck.read_bytes() == before


def test_inspect_reports_native_selectors_and_never_mutates_source(tmp_path: Path) -> None:
    deck = tmp_path / "inspect.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{
            "elements": [
                {"id": "first", "type": "shape", "x": 1, "y": 1, "width": 2, "height": 1, "text": "Alpha"},
                {"id": "second", "type": "text", "x": 4, "y": 1, "width": 2, "height": 1, "text": "Beta"},
            ],
            "animations": [{"type": "fade_in", "target": "first", "order": 0}],
            "transition": "fade",
        }],
    })
    presentation = Presentation(deck)
    presentation.slides[0].shapes[1].name = "first"
    presentation.save(deck)
    before = deck.read_bytes()

    inspected = inspect_presentation(deck)

    assert deck.read_bytes() == before
    assert inspected["operation"] == "inspect"
    assert inspected["presentation"]["slide_count"] == 1
    slide = inspected["slides"][0]
    assert slide["selector"] == {"slide_id": slide["slide_id"]}
    assert slide["part"] == "ppt/slides/slide1.xml"
    assert [shape["shape_id"] for shape in slide["shapes"]] == [2, 3]
    assert slide["shapes"][0]["selector"] == {
        "slide_id": slide["slide_id"], "shape_id": 2,
    }
    assert slide["shapes"][0]["text"] == "Alpha"
    assert slide["animations"][0]["target_shape_id"] == 2
    assert slide["transition"]["type"] == "fade"
    assert any(
        relationship["type"].endswith("/slideLayout")
        for relationship in slide["relationships"]
    )
    assert inspected["capabilities"]["stable_selector_edits"] is True
    assert inspected["preservation"]["source_mutated"] is False
    assert any(
        warning["code"] == "AMBIGUOUS_SHAPE_NAME"
        for warning in inspected["warnings"]
    )


def test_native_selectors_target_duplicate_names_and_survive_slide_moves(tmp_path: Path) -> None:
    deck = tmp_path / "native-selectors.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [
            {"elements": [
                {"id": "one", "type": "text", "x": 1, "y": 1, "width": 2, "height": 1, "text": "One"},
                {"id": "two", "type": "text", "x": 4, "y": 1, "width": 2, "height": 1, "text": "Two"},
            ]},
            {"elements": [
                {"id": "other", "type": "text", "x": 1, "y": 1, "width": 2, "height": 1, "text": "Other"},
            ]},
        ],
    })
    presentation = Presentation(deck)
    presentation.slides[0].shapes[0].name = "duplicate"
    presentation.slides[0].shapes[1].name = "duplicate"
    presentation.save(deck)
    inspected = inspect_presentation(deck)
    first_slide = inspected["slides"][0]
    first_shape, second_shape = first_slide["shapes"]

    execute_spec({
        "operation": "edit", "output_path": str(deck),
        "modifiers": [
            {
                "operation": "move_slide",
                "from_selector": first_slide["selector"],
                "to_index": 1,
            },
            {
                "operation": "update_element",
                "selector": second_shape["selector"],
                "changes": {"text": "Selected", "x": 5},
            },
            {
                "operation": "align_elements",
                "selectors": [first_shape["selector"], second_shape["selector"]],
                "alignment": "top",
            },
            {
                "operation": "set_transition",
                "selector": first_slide["selector"],
                "transition": "fade",
            },
        ],
    })
    reopened = Presentation(deck)
    assert [shape.text for shape in reopened.slides[1].shapes] == ["One", "Selected"]
    assert reopened.slides[1].shapes[1].left == Inches(5)
    assert reopened.slides[1].shapes[0].top == reopened.slides[1].shapes[1].top
    post_edit = inspect_presentation(deck)
    selected_slide = next(
        slide for slide in post_edit["slides"]
        if slide["slide_id"] == first_slide["slide_id"]
    )
    assert selected_slide["transition"]["type"] == "fade"


def test_ambiguous_legacy_shape_name_is_rejected_atomically(tmp_path: Path) -> None:
    deck = tmp_path / "ambiguous-name.pptx"
    execute_spec({
        "output_path": str(deck),
        "slides": [{"elements": [
            {"id": "a", "type": "text", "x": 1, "y": 1, "width": 2, "height": 1, "text": "A"},
            {"id": "b", "type": "text", "x": 4, "y": 1, "width": 2, "height": 1, "text": "B"},
        ]}],
    })
    presentation = Presentation(deck)
    presentation.slides[0].shapes[0].name = "same"
    presentation.slides[0].shapes[1].name = "same"
    presentation.save(deck)
    before = deck.read_bytes()

    with pytest.raises(ValueError, match="ambiguous.*native selector"):
        execute_spec({
            "operation": "edit", "output_path": str(deck),
            "modifiers": [{
                "operation": "update_element", "slide_index": 0,
                "element_id": "same", "changes": {"text": "Wrong"},
            }],
        })
    assert deck.read_bytes() == before


def test_json_cli_inspects_without_output_path(tmp_path: Path) -> None:
    deck = tmp_path / "inspect-cli.pptx"
    execute_spec({"output_path": str(deck), "slides": [{}]})
    before = deck.read_bytes()
    completed = subprocess.run(
        [sys.executable, "-m", "animated_pptx.cli"],
        input=json.dumps({"operation": "inspect", "input_path": str(deck)}),
        text=True,
        capture_output=True,
        check=False,
        env={
            **os.environ,
            "PYTHONPATH": str(
                Path(__file__).parents[1] / "src" / "agent" / "powerpoint"
            ),
        },
    )
    assert completed.returncode == 0, completed.stderr
    result = json.loads(completed.stdout)
    assert result["ok"] is True
    assert result["inspection"]["presentation"]["slide_count"] == 1
    assert deck.read_bytes() == before


def test_json_cli_generates_powerpoint(tmp_path: Path) -> None:
    output = tmp_path / "tool-generated.pptx"
    request = {
        "output_path": str(output),
        "slides": [{
            "elements": [{
                "id": "title", "type": "text", "x": 1, "y": 1,
                "width": 6, "height": 1, "text": "Specialized tool",
            }],
            "animations": [],
            "transition": None,
        }],
    }
    completed = subprocess.run(
        [sys.executable, "-m", "animated_pptx.cli"],
        input=json.dumps(request),
        text=True,
        capture_output=True,
        check=False,
        env={
            **os.environ,
            "PYTHONPATH": str(
                Path(__file__).parents[1] / "src" / "agent" / "powerpoint"
            ),
        },
    )
    assert completed.returncode == 0, completed.stderr
    result = json.loads(completed.stdout)
    assert result == {"ok": True, "path": str(output), "slides": 1}
    assert Presentation(output).slides[0].shapes[0].text == "Specialized tool"
